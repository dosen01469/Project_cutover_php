from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR


PRIMARY = RGBColor(31, 78, 121)
ACCENT = RGBColor(230, 126, 34)
LIGHT = RGBColor(242, 244, 247)
DARK = RGBColor(55, 65, 81)
WHITE = RGBColor(255, 255, 255)
RED = RGBColor(192, 57, 43)
GREEN = RGBColor(39, 174, 96)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_title(slide, title, subtitle=""):
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.1), Inches(1.0))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Calibri"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = PRIMARY

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(12.0), Inches(0.6))
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = "Calibri"
        r2.font.size = Pt(12)
        r2.font.color.rgb = DARK

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.65), Inches(2.4), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.color.rgb = ACCENT


def add_footer(slide, text="Project Cutover PHP Migration"):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(12.0), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(9)
    r.font.color.rgb = DARK


def add_bullets(slide, items, left=0.9, top=1.8, width=5.5, height=4.8, font_size=18, color=DARK):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = item
        p.level = 0
        p.font.name = "Calibri"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.bullet = True


def add_label_box(slide, text, left, top, width, height, fill_color, text_color=WHITE, font_size=16):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(font_size)
    r.font.bold = True
    r.font.color.rgb = text_color
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=ACCENT):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(2.5)
    line.line.end_arrowhead = True
    return line


# Slide 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Rencana Cutover Implementasi New Application", "Migrasi Sales Force, Distribution Core, dan Warehouse")
add_label_box(slide, "Sales Force\n10.100.10.7", 0.8, 2.2, 2.0, 1.0, PRIMARY)
add_label_box(slide, "Distribution Core\n10.100.10.28", 3.1, 2.2, 2.2, 1.0, PRIMARY)
add_label_box(slide, "Warehouse\n10.100.10.28", 5.6, 2.2, 2.0, 1.0, PRIMARY)
add_arrow(slide, 7.9, 2.7, 9.0, 2.7)
add_label_box(slide, "New Integrated App\n10.100.10.11", 9.1, 2.1, 3.0, 1.2, ACCENT)
add_label_box(slide, "DB 8.0\n10.100.10.12", 10.0, 3.7, 1.6, 0.9, DARK)
add_footer(slide)

# Slide 2
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Latar Belakang Perubahan")
add_bullets(slide, [
    "3 aplikasi besar masih berjalan terpisah",
    "Platform lama: PHP 5.7 dan MySQL 5.6",
    "Infrastruktur dan database tersebar di beberapa server/IP",
    "Aplikasi baru telah digabung dalam 1 environment",
    "Platform baru: PHP 7.4 dan MySQL 8.0",
    "Testing dan perbaikan modul telah selesai",
], left=0.9, top=1.9, width=5.8)
add_label_box(slide, "CURRENT", 7.7, 2.0, 2.0, 0.8, PRIMARY)
add_label_box(slide, "TARGET", 10.0, 2.0, 2.0, 0.8, ACCENT)
add_arrow(slide, 9.8, 2.4, 9.95, 2.4)
add_footer(slide)

# Slide 3
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Current State vs Target State")
add_label_box(slide, "CURRENT STATE", 0.9, 1.9, 5.3, 0.7, PRIMARY)
add_label_box(slide, "TARGET STATE", 7.0, 1.9, 5.3, 0.7, ACCENT)
add_bullets(slide, [
    "Sales Force App — 10.100.10.7",
    "Distribution Core — 10.100.10.28",
    "Warehouse Activity — 10.100.10.28",
    "Database lama tersebar",
    "PHP 5.7 / MySQL 5.6",
], left=1.0, top=2.9, width=5.0, font_size=17)
add_bullets(slide, [
    "New App terintegrasi",
    "App Server — 10.100.10.11",
    "DB Server — 10.100.10.12",
    "PHP 7.4 / MySQL 8.0",
    "1 environment terpusat",
], left=7.1, top=2.9, width=5.0, font_size=17)
add_footer(slide)

# Slide 4
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Tujuan Implementasi")
add_label_box(slide, "Integrasi", 1.0, 2.2, 2.2, 1.2, PRIMARY)
add_label_box(slide, "Stabilitas", 3.6, 2.2, 2.2, 1.2, ACCENT)
add_label_box(slide, "Upgrade Platform", 6.2, 2.2, 2.4, 1.2, PRIMARY)
add_label_box(slide, "Efisiensi Operasional", 9.0, 2.2, 3.0, 1.2, ACCENT)
add_bullets(slide, [
    "1 environment untuk 3 aplikasi",
    "Maintainability lebih baik",
    "Platform lebih update dan stabil",
    "Operasional lebih efisien dan terkontrol",
], left=1.2, top=4.2, width=10.5, font_size=18)
add_footer(slide)

# Slide 5
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Metode Cutover")
steps = [
    "Freeze",
    "Backup",
    "Migrasi",
    "Validasi",
    "Switch",
    "Monitor",
    "Rollback",
]
left = 0.7
for i, step in enumerate(steps):
    color = PRIMARY if i % 2 == 0 else ACCENT
    add_label_box(slide, step, left + i * 1.75, 2.8, 1.45, 0.9, color, font_size=14)
    if i < len(steps) - 1:
        add_arrow(slide, left + i * 1.75 + 1.45, 3.25, left + (i + 1) * 1.75, 3.25)
add_bullets(slide, [
    "Controlled Cutover / Big Bang dengan Rollback Plan",
    "Perpindahan dilakukan terencana dan terkontrol",
    "Jalur kembali disiapkan bila terjadi kendala kritikal",
], left=1.0, top=4.6, width=11.0, font_size=18)
add_footer(slide)

# Slide 6
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Tahapan Implementasi")
phases = [
    ("1. Preparation", PRIMARY, ["Final review", "Freeze change", "Backup readiness", "Communication"]),
    ("2. Cutover Execution", ACCENT, ["Stop transaksi", "Backup final", "Final migration", "Aktivasi New App"]),
    ("3. Validation & Go-Live", PRIMARY, ["Smoke test", "Key user validation", "Go / No-Go", "Switch access"]),
    ("4. Hypercare", ACCENT, ["Monitoring", "Issue handling", "Stabilization"]),
]
start_x = 0.6
for idx, (title, color, bullets) in enumerate(phases):
    x = start_x + idx * 3.15
    add_label_box(slide, title, x, 2.0, 2.7, 0.8, color, font_size=14)
    add_bullets(slide, bullets, left=x + 0.05, top=3.0, width=2.5, height=2.4, font_size=14)
add_footer(slide)

# Slide 7
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Alur Hari H Cutover")
flow_steps = [
    "Stop\nTransaksi",
    "Final\nBackup",
    "Migrasi\nFinal",
    "Validasi\nData & App",
    "Soft Opening\nKey User",
    "Go-Live\nFull User",
]
y = 1.8
for i, step in enumerate(flow_steps):
    color = PRIMARY if i % 2 == 0 else ACCENT
    add_label_box(slide, step, 4.8, y + i * 0.8, 3.0, 0.55, color, font_size=15)
    if i < len(flow_steps) - 1:
        add_arrow(slide, 6.3, y + i * 0.8 + 0.55, 6.3, y + (i + 1) * 0.8)
add_footer(slide)

# Slide 8
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Risiko Utama dan Mitigasi")
add_label_box(slide, "RISIKO", 0.9, 2.0, 4.7, 0.7, PRIMARY)
add_label_box(slide, "MITIGASI", 6.0, 2.0, 6.0, 0.7, ACCENT)
risks = [
    "Data tidak sesuai",
    "Transaksi inti gagal",
    "Performa lambat",
    "Integrasi tidak berjalan",
    "Gangguan saat live",
]
mitigations = [
    "Validasi dan rekonsiliasi data",
    "Smoke test dan key user test",
    "Monitoring server dan database",
    "Validasi interface sebelum go-live",
    "Rollback plan siap",
]
for i, (rsk, mit) in enumerate(zip(risks, mitigations)):
    top = 2.9 + i * 0.7
    add_label_box(slide, rsk, 0.9, top, 4.7, 0.5, LIGHT, text_color=DARK, font_size=14)
    add_label_box(slide, mit, 6.0, top, 6.0, 0.5, LIGHT, text_color=DARK, font_size=14)
add_footer(slide)

# Slide 9
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Backup Plan / Rollback Plan")
add_label_box(slide, "NEW APP", 8.6, 2.1, 2.4, 1.0, ACCENT)
add_label_box(slide, "OLD SYSTEM", 2.2, 2.1, 2.8, 1.0, PRIMARY)
add_arrow(slide, 8.5, 2.6, 5.3, 2.6, RED)
add_bullets(slide, [
    "Old system tidak langsung dimatikan",
    "Backup final dilakukan sebelum switch",
    "Jika ada kendala kritikal, akses dikembalikan ke old system",
    "Operasional bisnis tetap menjadi prioritas utama",
], left=1.0, top=4.1, width=11.0, font_size=18)
add_footer(slide)

# Slide 10
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Kesimpulan dan Permohonan Persetujuan")
add_bullets(slide, [
    "Aplikasi baru sudah siap secara fungsional",
    "Cutover akan dilakukan dengan proses terkontrol",
    "Validasi, monitoring, dan rollback plan telah disiapkan",
    "Implementasi dirancang untuk meminimalkan risiko operasional",
], left=1.0, top=2.0, width=10.8, font_size=20)
add_label_box(slide, "Permohonan Persetujuan Manajemen", 2.7, 5.0, 7.2, 0.9, GREEN, font_size=20)
add_footer(slide)


output_file = "cutover_implementation_presentation.pptx"
prs.save(output_file)
print(f"Presentation created: {output_file}")
